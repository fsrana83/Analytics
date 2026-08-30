# -*- coding: utf-8 -*-
"""
Smart Analytics
================
Executive analytics platform for an insurance company in Oman.
Single-file production-ready Streamlit application.

Run with:  streamlit run app.py

See the accompanying SETUP_GUIDE.md for installation, authentication and
Gemini API configuration instructions.
"""

import io
import json
import os
import hashlib
import datetime as dt
from typing import Optional

import numpy as np
import pandas as pd
import streamlit as st
import plotly.express as px
import plotly.graph_objects as go

# =============================================================================
# CONSTANTS / BRANDING
# =============================================================================

APP_NAME = "Smart Analytics"
CURRENCY_SYMBOL = "﷼"          # Official CBO Rial Omani symbol (see cbo.gov.om/omrsymbol)
CURRENCY_CODE = "OMR"
GEMINI_MODEL = "gemini-3.7-flash"
DATA_DIR = "smart_analytics_data"
HIERARCHY_FILE = os.path.join(DATA_DIR, "hierarchy.json")
CONFIG_FILE = os.path.join(DATA_DIR, "config.json")
USERS_FILE = os.path.join(DATA_DIR, "users.json")
DATA_FILE = os.path.join(DATA_DIR, "transactions.csv")

REQUIRED_COLUMNS = [
    "month", "year", "group_of_product", "product", "sub_product",
    "line_of_business", "branch", "channel", "agent",
    "gross_written_premium", "net_written_premium", "earned_premium",
    "claims_reported", "claims_paid", "outstanding_claims",
    "commissions", "brokerage_costs", "direct_costs", "indirect_costs",
    "reinsurance_costs", "receivables", "budget", "actual", "forecast",
]

# KPI definitions used throughout the app (glossary expanders + Management Report)
KPI_DEFINITIONS = {
    "Loss Ratio": "Claims paid divided by earned premium. Measures underwriting loss "
                   "experience — the share of premium consumed by claims.",
    "Expense Ratio": "Commissions, brokerage, direct and indirect costs combined, divided "
                      "by earned premium. Measures operating cost efficiency.",
    "Combined Ratio": "Loss Ratio plus Expense Ratio. Below 100% indicates an underwriting "
                       "profit; above 100% indicates an underwriting loss.",
    "Claims Ratio": "Claims reported divided by gross written premium. Indicates claims "
                     "frequency/severity relative to production volume.",
    "Receivables Turnover": "Gross written premium divided by receivables. Higher values "
                             "indicate faster premium collection.",
    "Collection Ratio": "Percentage of gross written premium already collected (i.e. not "
                         "sitting in receivables).",
    "Reinsurance Cost Ratio": "Reinsurance costs divided by gross written premium. Tracks "
                               "the cost of ceded risk relative to production.",
    "Profit Margin": "Net underwriting profitability divided by earned premium.",
    "Premium Growth": "Period-over-period percentage change in gross written premium.",
    "Budget Variance": "Difference between actual and budgeted production, expressed in "
                        "absolute value and as a percentage of budget.",
}

os.makedirs(DATA_DIR, exist_ok=True)


# =============================================================================
# NUMBER / CURRENCY FORMATTING (all figures rounded to 0 decimals)
# =============================================================================

def fmt_omr(value: float, decimals: int = 0) -> str:
    """Format a number as Rial Omani using the official CBO symbol prefix."""
    if value is None or (isinstance(value, float) and np.isnan(value)):
        return f"{CURRENCY_SYMBOL} -"
    sign = "-" if value < 0 else ""
    return f"{sign}{CURRENCY_SYMBOL} {abs(value):,.{decimals}f}"


def fmt_pct(value: float, decimals: int = 0) -> str:
    if value is None or (isinstance(value, float) and np.isnan(value)):
        return "-"
    return f"{value:,.{decimals}f}%"


def fmt_num(value: float, decimals: int = 0) -> str:
    if value is None or (isinstance(value, float) and np.isnan(value)):
        return "-"
    return f"{value:,.{decimals}f}"


# =============================================================================
# PERSISTENCE HELPERS
# =============================================================================

def _load_json(path, default):
    if os.path.exists(path):
        try:
            with open(path, "r", encoding="utf-8") as f:
                return json.load(f)
        except Exception:
            return default
    return default


def _save_json(path, data):
    with open(path, "w", encoding="utf-8") as f:
        json.dump(data, f, indent=2, default=str)


# ---- Hierarchy -------------------------------------------------------------
# Each entry's group_of_product doubles as the Line of Business, so demo
# data and filters cover every LOB consistently end-to-end.

def load_hierarchy() -> pd.DataFrame:
    default = [
        {"group_of_product": "Motor Insurance", "product": "Motor Third Party Cover", "sub_product": "Individual Vehicle"},
        {"group_of_product": "Motor Insurance", "product": "Motor Third Party Cover", "sub_product": "Fleet"},
        {"group_of_product": "Motor Insurance", "product": "Motor Third Party Cover", "sub_product": "Trade Plate"},
        {"group_of_product": "Motor Insurance", "product": "Motor Comprehensive Cover", "sub_product": "Individual Vehicle"},
        {"group_of_product": "Motor Insurance", "product": "Motor Extended Warranty Cover", "sub_product": "Standard"},
        {"group_of_product": "Property Insurance", "product": "Fire & Perils", "sub_product": "Commercial"},
        {"group_of_product": "Property Insurance", "product": "Fire & Perils", "sub_product": "Residential"},
        {"group_of_product": "Property Insurance", "product": "Property All Risks", "sub_product": "Commercial"},
        {"group_of_product": "Marine Insurance", "product": "Marine Cargo", "sub_product": "Import/Export"},
        {"group_of_product": "Marine Insurance", "product": "Marine Hull", "sub_product": "Commercial Vessels"},
        {"group_of_product": "Medical Insurance", "product": "Group Medical", "sub_product": "Corporate"},
        {"group_of_product": "Medical Insurance", "product": "Individual Medical", "sub_product": "Standard"},
        {"group_of_product": "Engineering Insurance", "product": "Contractors All Risks", "sub_product": "Construction Projects"},
        {"group_of_product": "Engineering Insurance", "product": "Machinery Breakdown", "sub_product": "Industrial"},
        {"group_of_product": "Liability Insurance", "product": "General Liability", "sub_product": "Commercial"},
        {"group_of_product": "Liability Insurance", "product": "Professional Indemnity", "sub_product": "Corporate"},
        {"group_of_product": "Travel Insurance", "product": "Individual Travel", "sub_product": "Standard"},
        {"group_of_product": "Travel Insurance", "product": "Group Travel", "sub_product": "Corporate"},
        {"group_of_product": "Life Insurance (Group)", "product": "Group Life", "sub_product": "Corporate"},
    ]
    rows = _load_json(HIERARCHY_FILE, default)
    return pd.DataFrame(rows)


def save_hierarchy(df: pd.DataFrame) -> None:
    _save_json(HIERARCHY_FILE, df.to_dict(orient="records"))


def validate_hierarchy(df: pd.DataFrame) -> list:
    """Return a list of validation error strings (empty = valid)."""
    errors = []
    required = {"group_of_product", "product", "sub_product"}
    if not required.issubset(set(df.columns)):
        errors.append(f"Hierarchy must contain columns: {sorted(required)}")
        return errors
    if df[["group_of_product", "product", "sub_product"]].isnull().any().any():
        errors.append("Hierarchy rows cannot contain empty Group/Product/Sub-product values.")
    dupes = df.duplicated(subset=["group_of_product", "product", "sub_product"])
    if dupes.any():
        errors.append(f"{dupes.sum()} duplicate hierarchy row(s) found.")
    return errors


# ---- Company configuration --------------------------------------------------

def load_config() -> dict:
    default = {
        "company_name": "A Demo Insurance",
        "currency_label": f"{CURRENCY_CODE} ({CURRENCY_SYMBOL})",
        "fiscal_year_start_month": 1,
        "gemini_api_key": "",
        "gemini_model": GEMINI_MODEL,
        "outlier_method": "Z-Score",
        "outlier_threshold": 3.0,
        "forecast_horizon_months": 12,
    }
    cfg = _load_json(CONFIG_FILE, default)
    for k, v in default.items():
        cfg.setdefault(k, v)
    return cfg


def save_config(cfg: dict) -> None:
    _save_json(CONFIG_FILE, cfg)


# ---- Users / auth ------------------------------------------------------------

def _hash_pw(password: str) -> str:
    return hashlib.sha256(password.encode("utf-8")).hexdigest()


def load_users() -> dict:
    default = {
        "admin": {"name": "System Administrator", "role": "Admin",
                   "password_hash": _hash_pw("admin123"), "email": "admin@smartanalytics.local"},
        "ceo": {"name": "Chief Executive Officer", "role": "CEO",
                "password_hash": _hash_pw("ceo123"), "email": "ceo@smartanalytics.local"},
        "coo": {"name": "Chief Operating Officer", "role": "COO",
                "password_hash": _hash_pw("coo123"), "email": "coo@smartanalytics.local"},
        "cfo": {"name": "Chief Financial Officer", "role": "CFO",
                "password_hash": _hash_pw("cfo123"), "email": "cfo@smartanalytics.local"},
    }
    return _load_json(USERS_FILE, default)


def save_users(users: dict) -> None:
    _save_json(USERS_FILE, users)


def authenticate_user(username: str, password: str) -> Optional[dict]:
    """Return the user record on success, else None.

    NOTE: This is a lightweight, dependency-free reference implementation
    (SHA-256 hashed local credential store). For a production deployment,
    swap this out for Streamlit's native OIDC auth (st.login) or the
    `streamlit-authenticator` package with cookies + bcrypt hashing —
    see SETUP_GUIDE.md.
    """
    users = load_users()
    user = users.get(username)
    if user and user["password_hash"] == _hash_pw(password):
        return {"username": username, **user}
    return None


def send_email_support(to_email: str, subject: str, body: str) -> bool:
    """Placeholder for password-recovery / support email delivery.

    Wire this to SMTP (smtplib), SendGrid, or AWS SES in production.
    Returns True on (simulated) success so the UI flow can be tested
    without real credentials configured.
    """
    st.session_state.setdefault("_email_outbox", [])
    st.session_state["_email_outbox"].append(
        {"to": to_email, "subject": subject, "body": body, "sent_at": str(dt.datetime.now())}
    )
    return True


# =============================================================================
# DATA LOAD / TEMPLATES / VALIDATION
# =============================================================================

def empty_dataset() -> pd.DataFrame:
    return pd.DataFrame(columns=REQUIRED_COLUMNS)


def load_data() -> pd.DataFrame:
    if os.path.exists(DATA_FILE):
        df = pd.read_csv(DATA_FILE)
        for c in REQUIRED_COLUMNS:
            if c not in df.columns:
                df[c] = np.nan
        return df
    return empty_dataset()


def save_data(df: pd.DataFrame) -> None:
    df.to_csv(DATA_FILE, index=False)


def generate_sample_data(hierarchy: pd.DataFrame, n_months: int = 24, seed: int = 42) -> pd.DataFrame:
    """Generate realistic sample data for testing, covering every LOB in the hierarchy."""
    rng = np.random.default_rng(seed)
    branches = ["Muscat HQ", "Salalah", "Sohar", "Nizwa", "Sur"]
    channels = ["Direct", "Broker", "Bancassurance", "Digital"]
    agents = [f"Agent-{i:03d}" for i in range(1, 16)]

    today = dt.date.today().replace(day=1)
    months = [(today - dt.timedelta(days=30 * i)) for i in range(n_months)][::-1]

    rows = []
    if hierarchy.empty:
        hierarchy = load_hierarchy()

    for m in months:
        for _, hrow in hierarchy.iterrows():
            for _ in range(rng.integers(1, 3)):
                gwp = float(rng.normal(45000, 12000))
                gwp = max(gwp, 2000)
                nwp = gwp * rng.uniform(0.7, 0.95)
                earned = nwp * rng.uniform(0.8, 1.05)
                claims_reported = gwp * rng.uniform(0.35, 0.75)
                claims_paid = claims_reported * rng.uniform(0.6, 0.95)
                outstanding = max(claims_reported - claims_paid, 0)
                commissions = gwp * rng.uniform(0.05, 0.15)
                brokerage = gwp * rng.uniform(0.01, 0.05)
                direct_costs = gwp * rng.uniform(0.03, 0.08)
                indirect_costs = gwp * rng.uniform(0.02, 0.05)
                reinsurance = gwp * rng.uniform(0.05, 0.20)
                receivables = gwp * rng.uniform(0.10, 0.40)
                budget = gwp * rng.uniform(0.9, 1.1)
                actual = gwp
                forecast = gwp * rng.uniform(0.92, 1.08)

                # occasional injected outlier
                if rng.random() < 0.03:
                    gwp *= rng.choice([2.5, 0.2])
                    actual = gwp

                rows.append({
                    "month": m.month, "year": m.year,
                    "group_of_product": hrow["group_of_product"],
                    "product": hrow["product"],
                    "sub_product": hrow["sub_product"],
                    "line_of_business": hrow["group_of_product"],  # LOB == product group, covers all LOBs
                    "branch": rng.choice(branches),
                    "channel": rng.choice(channels),
                    "agent": rng.choice(agents),
                    "gross_written_premium": round(gwp, 0),
                    "net_written_premium": round(nwp, 0),
                    "earned_premium": round(earned, 0),
                    "claims_reported": round(claims_reported, 0),
                    "claims_paid": round(claims_paid, 0),
                    "outstanding_claims": round(outstanding, 0),
                    "commissions": round(commissions, 0),
                    "brokerage_costs": round(brokerage, 0),
                    "direct_costs": round(direct_costs, 0),
                    "indirect_costs": round(indirect_costs, 0),
                    "reinsurance_costs": round(reinsurance, 0),
                    "receivables": round(receivables, 0),
                    "budget": round(budget, 0),
                    "actual": round(actual, 0),
                    "forecast": round(forecast, 0),
                })
    return pd.DataFrame(rows)


def make_template(template_type: str) -> pd.DataFrame:
    """Return an empty template DataFrame with the correct columns for a given upload type."""
    common = ["month", "year", "group_of_product", "product", "sub_product",
              "line_of_business", "branch", "channel", "agent"]
    templates = {
        "Budget": common + ["budget"],
        "Actual": common + ["actual"],
        "Forecast": common + ["forecast"],
        "Claims": common + ["claims_reported", "claims_paid", "outstanding_claims"],
        "Production": common + ["gross_written_premium", "net_written_premium", "earned_premium"],
        "Expenses": common + ["commissions", "brokerage_costs", "direct_costs", "indirect_costs", "reinsurance_costs"],
    }
    cols = templates.get(template_type, REQUIRED_COLUMNS)
    return pd.DataFrame(columns=cols)


def validate_upload(df: pd.DataFrame, template_type: str) -> list:
    """Validate an uploaded file against the expected template columns. Returns list of error strings."""
    errors = []
    template = make_template(template_type)
    missing = [c for c in template.columns if c not in df.columns]
    if missing:
        errors.append(f"Missing required column(s) for '{template_type}': {missing}")
        return errors

    hierarchy = load_hierarchy()
    valid_combos = set(
        tuple(x) for x in hierarchy[["group_of_product", "product", "sub_product"]].values
    )
    if not hierarchy.empty:
        bad_rows = 0
        for _, r in df.iterrows():
            combo = (r.get("group_of_product"), r.get("product"), r.get("sub_product"))
            if combo not in valid_combos:
                bad_rows += 1
        if bad_rows:
            errors.append(
                f"{bad_rows} row(s) reference a Group/Product/Sub-product combination "
                f"not defined in Setup. Please update Setup or correct the file."
            )

    numeric_candidates = [c for c in df.columns if c not in
                           ["month", "year", "group_of_product", "product", "sub_product",
                            "line_of_business", "branch", "channel", "agent"]]
    for c in numeric_candidates:
        non_numeric = pd.to_numeric(df[c], errors="coerce").isna() & df[c].notna()
        if non_numeric.any():
            errors.append(f"Column '{c}' contains {non_numeric.sum()} non-numeric value(s).")

    return errors


def merge_uploaded_data(existing: pd.DataFrame, new: pd.DataFrame) -> pd.DataFrame:
    for c in REQUIRED_COLUMNS:
        if c not in new.columns:
            new[c] = np.nan
    merged = pd.concat([existing, new[REQUIRED_COLUMNS]], ignore_index=True)
    return merged


# =============================================================================
# KPI CALCULATIONS
# =============================================================================

def calculate_kpis(df: pd.DataFrame) -> pd.DataFrame:
    """Compute derived ratios/KPIs row-wise (safe against divide-by-zero)."""
    out = df.copy()

    def safe_div(a, b):
        return np.where((b == 0) | pd.isna(b), np.nan, a / b)

    out["loss_ratio"] = safe_div(out["claims_paid"], out["earned_premium"]) * 100
    out["expense_ratio"] = safe_div(
        out["commissions"] + out["brokerage_costs"] + out["direct_costs"] + out["indirect_costs"],
        out["earned_premium"],
    ) * 100
    out["combined_ratio"] = out["loss_ratio"].fillna(0) + out["expense_ratio"].fillna(0)
    out["claims_ratio"] = safe_div(out["claims_reported"], out["gross_written_premium"]) * 100
    out["receivables_turnover"] = safe_div(out["gross_written_premium"], out["receivables"])
    out["collection_ratio"] = safe_div(
        out["gross_written_premium"] - out["receivables"], out["gross_written_premium"]
    ) * 100
    out["reinsurance_cost_ratio"] = safe_div(out["reinsurance_costs"], out["gross_written_premium"]) * 100
    out["profitability"] = (
        out["earned_premium"] - out["claims_paid"] - out["commissions"] - out["brokerage_costs"]
        - out["direct_costs"] - out["indirect_costs"] - out["reinsurance_costs"]
    )
    out["profit_margin"] = safe_div(out["profitability"], out["earned_premium"]) * 100
    out["variance_abs"] = out["actual"] - out["budget"]
    out["variance_pct"] = safe_div(out["actual"] - out["budget"], out["budget"]) * 100
    return out


def period_growth(df: pd.DataFrame, value_col: str = "gross_written_premium") -> pd.DataFrame:
    """Monthly premium growth % across the full portfolio."""
    g = df.groupby(["year", "month"], as_index=False)[value_col].sum().sort_values(["year", "month"])
    g["period"] = pd.to_datetime(dict(year=g.year, month=g.month, day=1))
    g["growth_pct"] = g[value_col].pct_change() * 100
    return g


def kpi_summary_table(kpi: pd.DataFrame) -> pd.DataFrame:
    """A single consolidated KPI summary table used across Home/CFO/Management Report."""
    rows = [
        ("Gross Written Premium", fmt_omr(kpi["gross_written_premium"].sum())),
        ("Net Written Premium", fmt_omr(kpi["net_written_premium"].sum())),
        ("Earned Premium", fmt_omr(kpi["earned_premium"].sum())),
        ("Loss Ratio", fmt_pct(np.nanmean(kpi["loss_ratio"]))),
        ("Expense Ratio", fmt_pct(np.nanmean(kpi["expense_ratio"]))),
        ("Combined Ratio", fmt_pct(np.nanmean(kpi["combined_ratio"]))),
        ("Claims Ratio", fmt_pct(np.nanmean(kpi["claims_ratio"]))),
        ("Receivables Turnover", fmt_num(np.nanmean(kpi["receivables_turnover"]))),
        ("Collection Ratio", fmt_pct(np.nanmean(kpi["collection_ratio"]))),
        ("Reinsurance Cost Ratio", fmt_pct(np.nanmean(kpi["reinsurance_cost_ratio"]))),
        ("Profitability", fmt_omr(kpi["profitability"].sum())),
        ("Profit Margin", fmt_pct(np.nanmean(kpi["profit_margin"]))),
    ]
    return pd.DataFrame(rows, columns=["KPI", "Value"])


# =============================================================================
# OUTLIER DETECTION
# =============================================================================

def detect_outliers(df: pd.DataFrame, value_col: str = "gross_written_premium",
                     group_cols=None, method: str = "Z-Score", threshold: float = 3.0) -> pd.DataFrame:
    """Flag unusual spikes/drops within groups (branch/agent/product/etc).

    method: 'Z-Score', 'IQR', or 'Isolation Forest'
    """
    if group_cols is None:
        group_cols = ["branch"]
    work = df.copy()
    work["is_outlier"] = False
    work["outlier_score"] = np.nan

    if work.empty:
        return work

    if method == "Isolation Forest":
        try:
            from sklearn.ensemble import IsolationForest
            X = work[[value_col]].fillna(work[value_col].median())
            model = IsolationForest(contamination=0.05, random_state=42)
            preds = model.fit_predict(X)
            scores = model.score_samples(X)
            work["is_outlier"] = preds == -1
            work["outlier_score"] = -scores
            return work
        except ImportError:
            method = "Z-Score"  # graceful fallback if scikit-learn isn't installed

    for _, gdf in work.groupby(group_cols):
        vals = gdf[value_col]
        if method == "IQR":
            q1, q3 = vals.quantile(0.25), vals.quantile(0.75)
            iqr = q3 - q1
            lower, upper = q1 - 1.5 * iqr, q3 + 1.5 * iqr
            flags = (vals < lower) | (vals > upper)
            score = (vals - vals.median()).abs()
        else:  # Z-Score
            mu, sigma = vals.mean(), vals.std(ddof=0)
            z = (vals - mu) / sigma if sigma and not np.isnan(sigma) and sigma != 0 else vals * 0
            flags = z.abs() > threshold
            score = z.abs()
        work.loc[gdf.index, "is_outlier"] = flags
        work.loc[gdf.index, "outlier_score"] = score

    return work


# =============================================================================
# FORECASTING
# =============================================================================

def forecast_metrics(df: pd.DataFrame, value_col: str = "gross_written_premium",
                      horizon: int = 12) -> pd.DataFrame:
    """Project the next `horizon` months using the best available lightweight method.

    Tries statsmodels Holt-Winters exponential smoothing first (captures trend +
    seasonality); falls back to a simple linear-trend + seasonal-naive blend if
    statsmodels is unavailable or history is too short.
    """
    hist = df.groupby(["year", "month"], as_index=False)[value_col].sum().sort_values(["year", "month"])
    hist["period"] = pd.to_datetime(dict(year=hist.year, month=hist.month, day=1))
    hist = hist.set_index("period")[value_col].asfreq("MS").interpolate()

    if len(hist) < 4:
        return pd.DataFrame(columns=["period", "value", "type", "lower", "upper"])

    future_index = pd.date_range(hist.index[-1] + pd.offsets.MonthBegin(1), periods=horizon, freq="MS")

    try:
        from statsmodels.tsa.holtwinters import ExponentialSmoothing
        seasonal_periods = 12 if len(hist) >= 24 else None
        model = ExponentialSmoothing(
            hist, trend="add",
            seasonal="add" if seasonal_periods else None,
            seasonal_periods=seasonal_periods,
            initialization_method="estimated",
        ).fit()
        preds = model.forecast(horizon)
        resid_std = np.std(model.resid) if hasattr(model, "resid") else hist.std() * 0.1
        lower = preds - 1.96 * resid_std
        upper = preds + 1.96 * resid_std
    except Exception:
        # Linear trend fallback
        x = np.arange(len(hist))
        coeffs = np.polyfit(x, hist.values, 1)
        trend = np.poly1d(coeffs)
        future_x = np.arange(len(hist), len(hist) + horizon)
        preds = pd.Series(trend(future_x), index=future_index)
        resid_std = np.std(hist.values - trend(x))
        lower = preds - 1.96 * resid_std
        upper = preds + 1.96 * resid_std

    hist_df = pd.DataFrame({"period": hist.index, "value": hist.values, "type": "Historical",
                             "lower": hist.values, "upper": hist.values})
    fc_df = pd.DataFrame({"period": future_index, "value": preds.values, "type": "Forecast",
                           "lower": lower.values, "upper": upper.values})
    return pd.concat([hist_df, fc_df], ignore_index=True)


# =============================================================================
# VARIANCE ANALYSIS (Budget vs Actual vs Forecast)
# =============================================================================

def variance_analysis(df: pd.DataFrame) -> dict:
    """Compute variance summary + driver breakdown for budget vs actual vs forecast."""
    kpi = calculate_kpis(df)
    total_budget = kpi["budget"].sum()
    total_actual = kpi["actual"].sum()
    total_forecast = kpi["forecast"].sum()

    variance_abs = total_actual - total_budget
    variance_pct = (variance_abs / total_budget * 100) if total_budget else np.nan
    fc_variance_abs = total_actual - total_forecast
    fc_variance_pct = (fc_variance_abs / total_forecast * 100) if total_forecast else np.nan

    by_group = kpi.groupby("group_of_product", as_index=False).agg(
        budget=("budget", "sum"), actual=("actual", "sum"), forecast=("forecast", "sum")
    )
    by_group["variance_abs"] = by_group["actual"] - by_group["budget"]
    by_group["variance_pct"] = np.where(
        by_group["budget"] != 0, by_group["variance_abs"] / by_group["budget"] * 100, np.nan
    )
    by_group["direction"] = np.where(by_group["variance_abs"] >= 0, "Favorable", "Unfavorable")

    drivers = {
        "Volume/Production": kpi["gross_written_premium"].sum() - kpi["budget"].sum(),
        "Claims": -(kpi["claims_paid"].sum()),
        "Reinsurance": -(kpi["reinsurance_costs"].sum()),
        "Expenses": -(kpi["commissions"].sum() + kpi["brokerage_costs"].sum()
                       + kpi["direct_costs"].sum() + kpi["indirect_costs"].sum()),
    }

    return {
        "total_budget": total_budget, "total_actual": total_actual, "total_forecast": total_forecast,
        "variance_abs": variance_abs, "variance_pct": variance_pct,
        "forecast_variance_abs": fc_variance_abs, "forecast_variance_pct": fc_variance_pct,
        "by_group": by_group, "drivers": drivers,
    }


# =============================================================================
# AI INSIGHTS (Gemini)
# =============================================================================

def generate_ai_insights(kpi_summary: dict, outliers_count: int, variance: dict,
                          api_key: str = "", model_name: str = GEMINI_MODEL) -> str:
    """Generate executive commentary. Uses Gemini if an API key is configured,
    otherwise falls back to a deterministic rule-based summary so the app is
    fully usable without any external AI dependency."""

    if api_key:
        try:
            return _generate_ai_insights_gemini(kpi_summary, outliers_count, variance, api_key, model_name)
        except Exception as e:
            st.warning(f"Gemini call failed, falling back to rule-based insights ({e}).")

    lines = []
    vpct = variance.get("variance_pct", np.nan)
    if not np.isnan(vpct):
        direction = "ahead of" if vpct >= 0 else "behind"
        lines.append(
            f"- Actual performance is running **{abs(vpct):.0f}% {direction} budget** "
            f"({fmt_omr(variance['variance_abs'])})."
        )
    lr = kpi_summary.get("loss_ratio")
    if lr is not None and not np.isnan(lr):
        risk = "elevated" if lr > 65 else "within a healthy range"
        lines.append(f"- The portfolio loss ratio stands at **{lr:.0f}%**, which is {risk}.")
    cr = kpi_summary.get("combined_ratio")
    if cr is not None and not np.isnan(cr):
        profitability_note = "indicating underwriting losses" if cr > 100 else "indicating underwriting profitability"
        lines.append(f"- Combined ratio of **{cr:.0f}%** is {profitability_note}.")
    if outliers_count:
        lines.append(
            f"- **{outliers_count} exception(s)** were flagged by outlier detection and warrant "
            f"branch/agent-level review."
        )
    rec_ratio = kpi_summary.get("collection_ratio")
    if rec_ratio is not None and not np.isnan(rec_ratio):
        if rec_ratio < 70:
            lines.append("- Receivables collection is lagging — recommend prioritizing follow-up on aged debtors.")

    lines.append(
        "- **Recommended actions:** review flagged exceptions with branch heads, validate reinsurance "
        "cost movements against treaty terms, and monitor forecast accuracy monthly."
    )
    if not lines:
        return "No significant trends detected for the selected period."
    return "\n".join(lines)


def _generate_ai_insights_gemini(kpi_summary: dict, outliers_count: int, variance: dict,
                                  api_key: str, model_name: str = GEMINI_MODEL) -> str:
    """Call the Google Gemini API to generate management commentary.

    Requires: pip install google-generativeai
    Configure the API key via the Admin/Configuration screen or environment variable.
    Default model: gemini-3.7-flash (override in Admin/Configuration).
    """
    import google.generativeai as genai

    genai.configure(api_key=api_key)
    model = genai.GenerativeModel(model_name)

    prompt = f"""You are the AI analyst for an insurance company's executive dashboard ("{APP_NAME}").
Write a concise (4-6 bullet points) executive commentary in {CURRENCY_CODE} covering:
what changed, why it changed, key risks, and recommended management actions.
Round all figures to the nearest whole number.

KPI summary: {json.dumps(kpi_summary, default=str)}
Outliers flagged: {outliers_count}
Variance summary: {json.dumps({k: v for k, v in variance.items() if k not in ('by_group',)}, default=str)}
"""
    response = model.generate_content(prompt)
    return response.text


# =============================================================================
# CHART HELPERS FOR PDF (matplotlib, used only inside PDF generation)
# =============================================================================

def _mpl_line_chart(x_labels, series: dict, title: str, ylabel: str = ""):
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    fig, ax = plt.subplots(figsize=(7.2, 3.2), dpi=150)
    for label, ys in series.items():
        ax.plot(x_labels, ys, marker="o", markersize=2.5, linewidth=1.6, label=label)
    ax.set_title(title, fontsize=11, color="#1f4e79", fontweight="bold")
    ax.set_ylabel(ylabel, fontsize=8)
    ax.tick_params(axis="x", labelrotation=45, labelsize=6)
    ax.tick_params(axis="y", labelsize=7)
    if len(series) > 1:
        ax.legend(fontsize=7)
    ax.grid(alpha=0.25)
    fig.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="png")
    plt.close(fig)
    buf.seek(0)
    return buf


def _mpl_bar_chart(labels, values, title: str, ylabel: str = ""):
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    fig, ax = plt.subplots(figsize=(7.2, 3.2), dpi=150)
    colors = ["#2ca02c" if v >= 0 else "#d62728" for v in values]
    ax.bar(labels, values, color=colors)
    ax.set_title(title, fontsize=11, color="#1f4e79", fontweight="bold")
    ax.set_ylabel(ylabel, fontsize=8)
    ax.tick_params(axis="x", labelrotation=30, labelsize=7)
    ax.tick_params(axis="y", labelsize=7)
    ax.grid(alpha=0.25, axis="y")
    fig.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="png")
    plt.close(fig)
    buf.seek(0)
    return buf


# =============================================================================
# REPORT GENERATION (Excel / CSV / PDF / DOCX / PPTX / Management Report)
# =============================================================================

def export_excel(dfs: dict) -> bytes:
    """dfs: {sheet_name: dataframe}"""
    buffer = io.BytesIO()
    with pd.ExcelWriter(buffer, engine="openpyxl") as writer:
        for name, d in dfs.items():
            d.to_excel(writer, sheet_name=name[:31], index=False)
    return buffer.getvalue()


def export_csv(df: pd.DataFrame) -> bytes:
    return df.to_csv(index=False).encode("utf-8")


def generate_reports(report_type: str, df: pd.DataFrame, variance: dict, insights_text: str,
                      company_name: str) -> bytes:
    """Generate a formatted report. report_type in {'PDF','Excel','PPTX','DOCX'}. Returns file bytes."""

    if report_type == "Excel":
        return export_excel({
            "Summary": pd.DataFrame([{
                "Company": company_name, "Total Budget": variance["total_budget"],
                "Total Actual": variance["total_actual"], "Total Forecast": variance["total_forecast"],
                "Variance %": round(variance["variance_pct"]) if not np.isnan(variance["variance_pct"]) else None,
            }]),
            "By Group": variance["by_group"],
            "Raw Data": df,
        })

    if report_type == "DOCX":
        from docx import Document
        doc = Document()
        doc.add_heading(f"{company_name} — {APP_NAME} Executive Report", level=1)
        doc.add_paragraph(f"Generated: {dt.datetime.now():%Y-%m-%d %H:%M}")
        doc.add_heading("Financial Summary", level=2)
        table = doc.add_table(rows=1, cols=2)
        table.style = "Light Grid Accent 1"
        hdr = table.rows[0].cells
        hdr[0].text, hdr[1].text = "Metric", "Value"
        summary_rows = [
            ("Total Budget", fmt_omr(variance["total_budget"])),
            ("Total Actual", fmt_omr(variance["total_actual"])),
            ("Total Forecast", fmt_omr(variance["total_forecast"])),
            ("Variance vs Budget", f"{fmt_omr(variance['variance_abs'])} ({fmt_pct(variance['variance_pct'])})"),
        ]
        for k, v in summary_rows:
            row = table.add_row().cells
            row[0].text, row[1].text = k, v
        doc.add_heading("AI Executive Insights", level=2)
        for line in insights_text.split("\n"):
            doc.add_paragraph(line.replace("**", ""), style="List Bullet" if line.strip().startswith("-") else None)
        buffer = io.BytesIO()
        doc.save(buffer)
        return buffer.getvalue()

    if report_type == "PPTX":
        from pptx import Presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = f"{company_name}"
        slide.placeholders[1].text = f"{APP_NAME} — Executive Report\n{dt.datetime.now():%Y-%m-%d}"

        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        slide2.shapes.title.text = "Financial Summary"
        body = slide2.placeholders[1].text_frame
        body.text = f"Total Budget: {fmt_omr(variance['total_budget'])}"
        for line in [
            f"Total Actual: {fmt_omr(variance['total_actual'])}",
            f"Total Forecast: {fmt_omr(variance['total_forecast'])}",
            f"Variance vs Budget: {fmt_omr(variance['variance_abs'])} ({fmt_pct(variance['variance_pct'])})",
        ]:
            p = body.add_paragraph()
            p.text = line

        slide3 = prs.slides.add_slide(prs.slide_layouts[1])
        slide3.shapes.title.text = "AI Executive Insights"
        body3 = slide3.placeholders[1].text_frame
        clean_lines = [l.replace("**", "").lstrip("- ") for l in insights_text.split("\n") if l.strip()]
        if clean_lines:
            body3.text = clean_lines[0]
            for line in clean_lines[1:]:
                p = body3.add_paragraph()
                p.text = line
        buffer = io.BytesIO()
        prs.save(buffer)
        return buffer.getvalue()

    if report_type == "PDF":
        return _generate_portfolio_pdf(df, variance, insights_text, company_name)

    raise ValueError(f"Unsupported report type: {report_type}")


def _generate_portfolio_pdf(df: pd.DataFrame, variance: dict, insights_text: str, company_name: str) -> bytes:
    """Portfolio summary PDF — headline financials + AI Executive Insights."""
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import cm
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
    from reportlab.lib import colors
    from reportlab.lib.styles import getSampleStyleSheet

    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=A4)
    styles = getSampleStyleSheet()
    story = [
        Paragraph(f"{company_name}", styles["Title"]),
        Paragraph(f"{APP_NAME} — Portfolio Report", styles["Heading2"]),
        Paragraph(f"Generated: {dt.datetime.now():%Y-%m-%d %H:%M}", styles["Normal"]),
        Spacer(1, 0.5 * cm),
    ]
    data = [["Metric", "Value"],
            ["Total Budget", fmt_omr(variance["total_budget"])],
            ["Total Actual", fmt_omr(variance["total_actual"])],
            ["Total Forecast", fmt_omr(variance["total_forecast"])],
            ["Variance vs Budget", f"{fmt_omr(variance['variance_abs'])} ({fmt_pct(variance['variance_pct'])})"]]
    t = Table(data, colWidths=[8 * cm, 8 * cm])
    t.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1f4e79")),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
        ("FONTSIZE", (0, 0), (-1, -1), 9),
    ]))
    story.append(t)
    story.append(Spacer(1, 0.5 * cm))

    # --- AI Executive Insights (always included in the portfolio PDF) ---
    story.append(Paragraph("AI Executive Insights", styles["Heading2"]))
    for line in insights_text.split("\n"):
        clean = line.replace("**", "")
        story.append(Paragraph(clean, styles["Normal"]))
    doc.build(story)
    return buffer.getvalue()


def generate_management_report(df: pd.DataFrame, cfg: dict, insights_text: str) -> bytes:
    """Rich, fully formatted management report (PDF): KPI dashboard with
    explanations, variance analysis, production forecast (chart + table),
    and AI executive insights. This is the comprehensive board-level
    deliverable, distinct from the shorter Portfolio Report."""
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import cm
    from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle,
                                     Image, PageBreak, HRFlowable)
    from reportlab.lib import colors
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.enums import TA_CENTER

    company_name = cfg["company_name"]
    kpi = calculate_kpis(df)
    variance = variance_analysis(df)
    outliers = detect_outliers(kpi, method=cfg.get("outlier_method", "Z-Score"),
                                threshold=cfg.get("outlier_threshold", 3.0))
    outliers_count = int(outliers["is_outlier"].sum())

    styles = getSampleStyleSheet()
    title_style = ParagraphStyle("CoverTitle", parent=styles["Title"], fontSize=26, alignment=TA_CENTER,
                                  textColor=colors.HexColor("#1f4e79"))
    subtitle_style = ParagraphStyle("CoverSub", parent=styles["Normal"], fontSize=13, alignment=TA_CENTER,
                                     textColor=colors.HexColor("#566073"))
    section_style = ParagraphStyle("Section", parent=styles["Heading1"], fontSize=15,
                                    textColor=colors.HexColor("#1f4e79"), spaceBefore=14, spaceAfter=8)
    body_style = styles["Normal"]

    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=A4, topMargin=2 * cm, bottomMargin=2 * cm)
    story = []

    # ---- Cover ----
    story.append(Spacer(1, 5 * cm))
    story.append(Paragraph(company_name, title_style))
    story.append(Spacer(1, 0.3 * cm))
    story.append(Paragraph(f"{APP_NAME} — Full Management Report", subtitle_style))
    story.append(Spacer(1, 0.2 * cm))
    period_min = f"{int(df['year'].min())}" if not df.empty else "-"
    period_max = f"{int(df['year'].max())}" if not df.empty else "-"
    story.append(Paragraph(f"Reporting Period: {period_min} – {period_max}", subtitle_style))
    story.append(Paragraph(f"Generated: {dt.datetime.now():%Y-%m-%d %H:%M}", subtitle_style))
    story.append(PageBreak())

    # ---- Executive Summary (AI Insights) ----
    story.append(Paragraph("Executive Summary", section_style))
    story.append(HRFlowable(width="100%", color=colors.HexColor("#1f4e79"), thickness=1))
    story.append(Spacer(1, 0.3 * cm))
    for line in insights_text.split("\n"):
        story.append(Paragraph(line.replace("**", ""), body_style))
    story.append(Spacer(1, 0.4 * cm))

    # ---- KPI Dashboard with explanations ----
    story.append(Paragraph("KPI Dashboard", section_style))
    story.append(HRFlowable(width="100%", color=colors.HexColor("#1f4e79"), thickness=1))
    story.append(Spacer(1, 0.3 * cm))

    kpi_values = {
        "Gross Written Premium": fmt_omr(kpi["gross_written_premium"].sum()),
        "Loss Ratio": fmt_pct(np.nanmean(kpi["loss_ratio"])),
        "Expense Ratio": fmt_pct(np.nanmean(kpi["expense_ratio"])),
        "Combined Ratio": fmt_pct(np.nanmean(kpi["combined_ratio"])),
        "Claims Ratio": fmt_pct(np.nanmean(kpi["claims_ratio"])),
        "Receivables Turnover": fmt_num(np.nanmean(kpi["receivables_turnover"])),
        "Collection Ratio": fmt_pct(np.nanmean(kpi["collection_ratio"])),
        "Reinsurance Cost Ratio": fmt_pct(np.nanmean(kpi["reinsurance_cost_ratio"])),
        "Profit Margin": fmt_pct(np.nanmean(kpi["profit_margin"])),
        "Budget Variance": f"{fmt_omr(variance['variance_abs'])} ({fmt_pct(variance['variance_pct'])})",
    }
    kpi_table_data = [["KPI", "Value", "Definition"]]
    for k, v in kpi_values.items():
        definition = KPI_DEFINITIONS.get(k, "")
        kpi_table_data.append([k, v, Paragraph(definition, ParagraphStyle("small", fontSize=7.5, leading=9))])
    kt = Table(kpi_table_data, colWidths=[4 * cm, 3 * cm, 9 * cm], repeatRows=1)
    kt.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1f4e79")),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("FONTSIZE", (0, 0), (-1, 0), 8),
        ("FONTSIZE", (0, 1), (1, -1), 8),
        ("GRID", (0, 0), (-1, -1), 0.4, colors.grey),
        ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#F7F9FC")]),
    ]))
    story.append(kt)
    story.append(Paragraph(f"Outlier exceptions flagged this period: {outliers_count}", body_style))
    story.append(PageBreak())

    # ---- Variance Analysis ----
    story.append(Paragraph("Variance Analysis — Budget vs Actual vs Forecast", section_style))
    story.append(HRFlowable(width="100%", color=colors.HexColor("#1f4e79"), thickness=1))
    story.append(Spacer(1, 0.3 * cm))
    story.append(Paragraph(
        f"Total Budget: {fmt_omr(variance['total_budget'])} · Total Actual: {fmt_omr(variance['total_actual'])} "
        f"· Total Forecast: {fmt_omr(variance['total_forecast'])}", body_style))
    story.append(Paragraph(
        f"Overall variance vs budget: {fmt_omr(variance['variance_abs'])} ({fmt_pct(variance['variance_pct'])}), "
        f"variance vs forecast: {fmt_omr(variance['forecast_variance_abs'])} "
        f"({fmt_pct(variance['forecast_variance_pct'])}).", body_style))
    story.append(Spacer(1, 0.25 * cm))

    by_group = variance["by_group"]
    var_chart_buf = _mpl_bar_chart(
        by_group["group_of_product"].tolist(), by_group["variance_abs"].round(0).tolist(),
        "Variance vs Budget by Product Group (OMR)")
    story.append(Image(var_chart_buf, width=16 * cm, height=7 * cm))
    story.append(Spacer(1, 0.2 * cm))

    var_table_data = [["Group of Product", "Budget", "Actual", "Variance", "Variance %"]]
    for _, r in by_group.iterrows():
        var_table_data.append([r["group_of_product"], fmt_omr(r["budget"]), fmt_omr(r["actual"]),
                                fmt_omr(r["variance_abs"]), fmt_pct(r["variance_pct"])])
    vt = Table(var_table_data, colWidths=[5 * cm, 3.5 * cm, 3.5 * cm, 3.5 * cm, 2.5 * cm], repeatRows=1)
    vt.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1f4e79")),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("FONTSIZE", (0, 0), (-1, -1), 7.5),
        ("GRID", (0, 0), (-1, -1), 0.4, colors.grey),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#F7F9FC")]),
    ]))
    story.append(vt)
    story.append(Spacer(1, 0.3 * cm))

    story.append(Paragraph("Variance Drivers", styles["Heading3"]))
    driver_table_data = [["Driver", "Impact (OMR)"]]
    for k, v in variance["drivers"].items():
        driver_table_data.append([k, fmt_omr(v)])
    dt_table = Table(driver_table_data, colWidths=[7 * cm, 6 * cm])
    dt_table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1f4e79")),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("FONTSIZE", (0, 0), (-1, -1), 8),
        ("GRID", (0, 0), (-1, -1), 0.4, colors.grey),
    ]))
    story.append(dt_table)
    story.append(PageBreak())

    # ---- Production Forecast ----
    story.append(Paragraph("Production Forecast", section_style))
    story.append(HRFlowable(width="100%", color=colors.HexColor("#1f4e79"), thickness=1))
    story.append(Spacer(1, 0.3 * cm))
    horizon = cfg.get("forecast_horizon_months", 12)
    fdf = forecast_metrics(kpi, horizon=horizon)
    if fdf.empty:
        story.append(Paragraph("Not enough historical data to forecast (need at least 4 months).", body_style))
    else:
        hist = fdf[fdf["type"] == "Historical"]
        fut = fdf[fdf["type"] == "Forecast"]
        x_labels = [d.strftime("%b-%y") for d in fdf["period"]]
        series = {}
        combined_vals = list(hist["value"]) + [np.nan] * len(fut)
        forecast_vals = [np.nan] * len(hist) + list(fut["value"])
        series["Historical"] = combined_vals
        series["Forecast"] = forecast_vals
        fc_chart_buf = _mpl_line_chart(x_labels, series,
                                        f"Gross Written Premium — {horizon}-Month Forecast", "OMR")
        story.append(Image(fc_chart_buf, width=16 * cm, height=7 * cm))
        story.append(Spacer(1, 0.2 * cm))

        fc_table_data = [["Period", "Forecast (OMR)", "Lower CI", "Upper CI"]]
        for _, r in fut.iterrows():
            fc_table_data.append([r["period"].strftime("%b %Y"), fmt_omr(r["value"]),
                                   fmt_omr(r["lower"]), fmt_omr(r["upper"])])
        ft = Table(fc_table_data, colWidths=[4 * cm, 4 * cm, 4 * cm, 4 * cm], repeatRows=1)
        ft.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1f4e79")),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("FONTSIZE", (0, 0), (-1, -1), 7.5),
            ("GRID", (0, 0), (-1, -1), 0.4, colors.grey),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#F7F9FC")]),
        ]))
        story.append(ft)

    story.append(Spacer(1, 0.5 * cm))
    story.append(Paragraph(
        f"Report generated by {APP_NAME} on behalf of {company_name}. All figures in "
        f"{CURRENCY_CODE} ({CURRENCY_SYMBOL}), rounded to the nearest whole number.",
        ParagraphStyle("footer", fontSize=7.5, textColor=colors.grey)))

    doc.build(story)
    return buffer.getvalue()


# =============================================================================
# PER-ROLE VIEW EXPORT (CEO / COO / CFO — PDF & PPTX)
# =============================================================================

ROLE_VIEW_TITLES = {
    "CEO": "CEO View — Strategic Performance",
    "COO": "COO View — Operations & Production",
    "CFO": "CFO View — Financial Health",
}


def _kpi_dict_for_role(role: str, kpi: pd.DataFrame, var: dict) -> dict:
    if role == "CEO":
        return {
            "Total GWP": fmt_omr(kpi["gross_written_premium"].sum()),
            "Profitability": fmt_omr(kpi["profitability"].sum()),
            "Combined Ratio": fmt_pct(kpi["combined_ratio"].mean()),
            "Budget vs Actual": fmt_pct(var["variance_pct"]),
        }
    if role == "COO":
        return {
            "Claims Reported": fmt_omr(kpi["claims_reported"].sum()),
            "Claims Paid": fmt_omr(kpi["claims_paid"].sum()),
            "Outstanding Claims": fmt_omr(kpi["outstanding_claims"].sum()),
            "Claims Ratio": fmt_pct(kpi["claims_ratio"].mean()),
        }
    return {  # CFO
        "Loss Ratio": fmt_pct(kpi["loss_ratio"].mean()),
        "Expense Ratio": fmt_pct(kpi["expense_ratio"].mean()),
        "Receivables": fmt_omr(kpi["receivables"].sum()),
        "Collection Ratio": fmt_pct(kpi["collection_ratio"].mean()),
    }


def _charts_for_role(role: str, kpi: pd.DataFrame) -> list:
    """Return [(title, png_bytesio), ...] rendered with matplotlib for embedding in PDF/PPTX."""
    charts = []
    if role == "CEO":
        by_group = kpi.groupby("group_of_product", as_index=False)["profitability"].sum()
        charts.append(("Profitability by Line of Business", _mpl_bar_chart(
            by_group["group_of_product"].tolist(), by_group["profitability"].round(0).tolist(),
            "Profitability by Line of Business (OMR)")))
    elif role == "COO":
        by_branch = kpi.groupby("branch", as_index=False)["gross_written_premium"].sum()
        charts.append(("Production by Branch", _mpl_bar_chart(
            by_branch["branch"].tolist(), by_branch["gross_written_premium"].round(0).tolist(),
            "Production by Branch (OMR)")))
        by_channel = kpi.groupby("channel", as_index=False)["gross_written_premium"].sum()
        charts.append(("Production by Channel", _mpl_bar_chart(
            by_channel["channel"].tolist(), by_channel["gross_written_premium"].round(0).tolist(),
            "Production by Channel (OMR)")))
    else:  # CFO
        trend = kpi.groupby(["year", "month"], as_index=False).agg(
            loss_ratio=("loss_ratio", "mean"), expense_ratio=("expense_ratio", "mean"),
            combined_ratio=("combined_ratio", "mean")).sort_values(["year", "month"])
        x_labels = [f"{int(r.year)}-{int(r.month):02d}" for r in trend.itertuples()]
        series = {
            "Loss Ratio": trend["loss_ratio"].round(1).tolist(),
            "Expense Ratio": trend["expense_ratio"].round(1).tolist(),
            "Combined Ratio": trend["combined_ratio"].round(1).tolist(),
        }
        charts.append(("Ratio Trends", _mpl_line_chart(x_labels, series, "Ratio Trends", "%")))
    return charts


def generate_view_export(role: str, df: pd.DataFrame, cfg: dict, insights_text: str, fmt: str) -> bytes:
    """Export a single role view (CEO/COO/CFO) as a fully formatted PDF or PPTX,
    including its KPIs, charts, and the current AI Insights."""
    kpi = calculate_kpis(df)
    var = variance_analysis(df)
    kpi_dict = _kpi_dict_for_role(role, kpi, var)
    charts = _charts_for_role(role, kpi)
    company_name = cfg["company_name"]
    title = ROLE_VIEW_TITLES[role]

    if fmt == "PDF":
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.units import cm
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image
        from reportlab.lib import colors
        from reportlab.lib.styles import getSampleStyleSheet

        styles = getSampleStyleSheet()
        buffer = io.BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=A4)
        story = [
            Paragraph(company_name, styles["Title"]),
            Paragraph(f"{APP_NAME} — {title}", styles["Heading2"]),
            Paragraph(f"Generated: {dt.datetime.now():%Y-%m-%d %H:%M}", styles["Normal"]),
            Spacer(1, 0.4 * cm),
        ]
        kpi_table_data = [["KPI", "Value"]] + [[k, v] for k, v in kpi_dict.items()]
        t = Table(kpi_table_data, colWidths=[8 * cm, 8 * cm])
        t.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1f4e79")),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
            ("FONTSIZE", (0, 0), (-1, -1), 9),
        ]))
        story.append(t)
        story.append(Spacer(1, 0.4 * cm))
        for chart_title, buf in charts:
            story.append(Paragraph(chart_title, styles["Heading3"]))
            story.append(Image(buf, width=16 * cm, height=7 * cm))
            story.append(Spacer(1, 0.3 * cm))
        story.append(Paragraph("AI Insights", styles["Heading2"]))
        for line in insights_text.split("\n"):
            story.append(Paragraph(line.replace("**", ""), styles["Normal"]))
        doc.build(story)
        return buffer.getvalue()

    if fmt == "PPTX":
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = company_name
        slide.placeholders[1].text = f"{APP_NAME} — {title}\n{dt.datetime.now():%Y-%m-%d}"

        kpi_slide = prs.slides.add_slide(prs.slide_layouts[1])
        kpi_slide.shapes.title.text = "Key Performance Indicators"
        body = kpi_slide.placeholders[1].text_frame
        items = list(kpi_dict.items())
        body.text = f"{items[0][0]}: {items[0][1]}"
        for k, v in items[1:]:
            p = body.add_paragraph()
            p.text = f"{k}: {v}"

        for chart_title, buf in charts:
            cslide = prs.slides.add_slide(prs.slide_layouts[5])
            cslide.shapes.title.text = chart_title
            buf.seek(0)
            cslide.shapes.add_picture(buf, Inches(0.7), Inches(1.5), width=Inches(8.5))

        insight_slide = prs.slides.add_slide(prs.slide_layouts[1])
        insight_slide.shapes.title.text = "AI Insights"
        ibody = insight_slide.placeholders[1].text_frame
        clean_lines = [l.replace("**", "").lstrip("- ") for l in insights_text.split("\n") if l.strip()]
        if clean_lines:
            ibody.text = clean_lines[0]
            for line in clean_lines[1:]:
                p = ibody.add_paragraph()
                p.text = line

        buffer = io.BytesIO()
        prs.save(buffer)
        return buffer.getvalue()

    raise ValueError(f"Unsupported format: {fmt}")


def render_ai_insights_and_export(role: str, df: pd.DataFrame, cfg: dict):
    """Shared block rendered at the bottom of each CEO/COO/CFO view: AI Insights
    (cached per role) plus PDF/PPTX export of that formatted view."""
    kpi = calculate_kpis(df)
    var = variance_analysis(df)
    outliers = detect_outliers(kpi, method=cfg.get("outlier_method", "Z-Score"),
                                threshold=cfg.get("outlier_threshold", 3.0))
    outliers_count = int(outliers["is_outlier"].sum())
    kpi_summary = {
        "loss_ratio": kpi["loss_ratio"].mean(),
        "expense_ratio": kpi["expense_ratio"].mean(),
        "combined_ratio": kpi["combined_ratio"].mean(),
        "collection_ratio": kpi["collection_ratio"].mean(),
        "profitability": kpi["profitability"].sum(),
    }
    cache_key = f"_ai_insights_{role.lower()}"

    st.markdown("---")
    header_col, btn_col = st.columns([5, 1])
    header_col.markdown("### 🤖 AI Insights")
    if btn_col.button("🔄 Regenerate", key=f"regen_{role}"):
        st.session_state.pop(cache_key, None)

    if cache_key not in st.session_state:
        with st.spinner("Analyzing KPIs, outliers, and variances..."):
            st.session_state[cache_key] = generate_ai_insights(
                kpi_summary, outliers_count, var,
                api_key=cfg.get("gemini_api_key", ""),
                model_name=cfg.get("gemini_model", GEMINI_MODEL))
    st.markdown(st.session_state[cache_key])

    st.markdown("##### 📤 Export this view")
    ec1, ec2 = st.columns(2)
    with ec1:
        if st.button(f"Generate PDF", key=f"gen_pdf_{role}"):
            with st.spinner("Building formatted PDF..."):
                content = generate_view_export(role, df, cfg, st.session_state[cache_key], "PDF")
            st.download_button(f"⬇️ Download {role} View (PDF)", data=content,
                                file_name=f"smart_analytics_{role.lower()}_view.pdf", mime="application/pdf",
                                key=f"dl_pdf_{role}")
    with ec2:
        if st.button(f"Generate PPTX", key=f"gen_pptx_{role}"):
            with st.spinner("Building formatted PPTX..."):
                content = generate_view_export(role, df, cfg, st.session_state[cache_key], "PPTX")
            st.download_button(
                f"⬇️ Download {role} View (PPTX)", data=content,
                file_name=f"smart_analytics_{role.lower()}_view.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                key=f"dl_pptx_{role}")


# =============================================================================
# STREAMLIT APP CONFIG / STYLE
# =============================================================================

st.set_page_config(page_title=APP_NAME, page_icon="📊", layout="wide", initial_sidebar_state="expanded")

CUSTOM_CSS = """
<style>
    .main { background-color: #f7f9fc; }
    div[data-testid="stMetric"] {
        background-color: white; border: 1px solid #e6e9ef; border-radius: 10px;
        padding: 12px 14px; box-shadow: 0 1px 3px rgba(0,0,0,0.04); overflow: hidden;
    }
    div[data-testid="stMetricLabel"] { font-weight: 600; color: #566073; font-size: 0.8rem; }
    div[data-testid="stMetricValue"] {
        font-size: 1.35rem !important; line-height: 1.25 !important;
        white-space: normal !important; word-break: break-word;
    }
    div[data-testid="stMetricDelta"] { font-size: 0.78rem !important; }
    h1, h2, h3 { color: #1f4e79; }
    .stTabs [data-baseweb="tab"] { font-weight: 600; }
    .exec-badge {
        display:inline-block; padding:2px 10px; border-radius:12px;
        background:#1f4e79; color:white; font-size:0.75rem; font-weight:600;
    }
</style>
"""
st.markdown(CUSTOM_CSS, unsafe_allow_html=True)


# =============================================================================
# SESSION STATE INIT
# =============================================================================

def init_state():
    st.session_state.setdefault("authenticated", False)
    st.session_state.setdefault("user", None)
    st.session_state.setdefault("page", "Home Dashboard")


init_state()


# =============================================================================
# LOGIN SCREEN
# =============================================================================

def login_screen():
    st.markdown(f"## 📊 {APP_NAME}")
    st.caption("Executive Insurance Analytics — Sultanate of Oman")
    col1, col2 = st.columns([1, 1])
    with col1:
        with st.form("login_form"):
            st.subheader("Sign in")
            username = st.text_input("Username")
            password = st.text_input("Password", type="password")
            submitted = st.form_submit_button("Login", use_container_width=True)
            if submitted:
                user = authenticate_user(username, password)
                if user:
                    st.session_state["authenticated"] = True
                    st.session_state["user"] = user
                    st.rerun()
                else:
                    st.error("Invalid username or password.")
        with st.expander("Forgot password?"):
            recovery_email = st.text_input("Registered email", key="recovery_email")
            if st.button("Send recovery instructions"):
                if send_email_support(recovery_email, "Password Recovery",
                                       "A password reset was requested for your Smart Analytics account."):
                    st.success("Recovery instructions sent (simulated — configure SMTP in production).")
    with col2:
        st.info(
            "**Demo credentials**\n\n"
            "- `admin` / `admin123` — Admin\n"
            "- `ceo` / `ceo123` — CEO view\n"
            "- `coo` / `coo123` — COO view\n"
            "- `cfo` / `cfo123` — CFO view\n\n"
            "Change these immediately in Admin → User Management before production use."
        )


# =============================================================================
# SIDEBAR / FILTERS
# =============================================================================

def sidebar_nav(cfg: dict, hierarchy: pd.DataFrame):
    user = st.session_state["user"]
    st.sidebar.markdown(f"### 📊 {APP_NAME}")
    st.sidebar.markdown(f"**{cfg['company_name']}**")
    st.sidebar.markdown(f"👤 {user['name']}  \n`{user['role']}`")
    if st.sidebar.button("Logout", use_container_width=True):
        st.session_state["authenticated"] = False
        st.session_state["user"] = None
        st.rerun()
    st.sidebar.divider()

    pages = ["Home Dashboard", "Setup", "Data Upload", "CEO View", "COO View", "CFO View",
              "Outlier Detection", "Forecasting", "Budget vs Actual vs Forecast",
              "AI Insights", "Reports"]
    if user["role"] == "Admin":
        pages.append("Admin / Configuration")
    page = st.sidebar.radio("Navigate", pages, index=pages.index(st.session_state["page"])
                             if st.session_state["page"] in pages else 0)
    st.session_state["page"] = page

    st.sidebar.divider()
    st.sidebar.markdown("#### Filters")
    df = load_data()
    filters = {}
    if not df.empty:
        years = sorted(df["year"].dropna().unique().tolist())
        filters["year"] = st.sidebar.multiselect("Year", years, default=years)
        months = sorted(df["month"].dropna().unique().tolist())
        filters["month"] = st.sidebar.multiselect("Month", months, default=months)
        if not hierarchy.empty:
            groups = sorted(hierarchy["group_of_product"].unique().tolist())
            filters["group_of_product"] = st.sidebar.multiselect("Line of Business / Group", groups, default=groups)
        for col, label in [("branch", "Branch"), ("channel", "Channel"), ("agent", "Agent")]:
            vals = sorted(df[col].dropna().unique().tolist())
            filters[col] = st.sidebar.multiselect(label, vals, default=vals)
    return page, filters


def apply_filters(df: pd.DataFrame, filters: dict) -> pd.DataFrame:
    out = df.copy()
    for col, selected in filters.items():
        if selected and col in out.columns:
            out = out[out[col].isin(selected)]
    return out


def kpi_glossary_expander():
    with st.expander("📖 KPI Glossary — what each metric means"):
        for k, v in KPI_DEFINITIONS.items():
            st.markdown(f"**{k}** — {v}")


# =============================================================================
# PAGE: HOME DASHBOARD
# =============================================================================

def page_home(df: pd.DataFrame, cfg: dict):
    st.title("🏠 Home Dashboard")
    st.caption(f"{cfg['company_name']} — All figures in {cfg['currency_label']}, rounded to the nearest whole number")

    if df.empty:
        st.warning("No data available yet. Go to **Data Upload** to load data, or generate sample data in **Admin / Configuration**.")
        return

    kpi = calculate_kpis(df)
    gwp = kpi["gross_written_premium"].sum()
    growth = period_growth(kpi)
    growth_last = growth["growth_pct"].iloc[-1] if len(growth) > 1 else 0

    c1, c2, c3, c4, c5 = st.columns(5)
    c1.metric("Gross Written Premium", fmt_omr(gwp), f"{growth_last:+.0f}% MoM")
    c2.metric("Loss Ratio", fmt_pct(kpi["claims_paid"].sum() / kpi["earned_premium"].sum() * 100
                                     if kpi["earned_premium"].sum() else np.nan))
    c3.metric("Combined Ratio", fmt_pct(
        (kpi["claims_paid"].sum() + kpi["commissions"].sum() + kpi["brokerage_costs"].sum()
         + kpi["direct_costs"].sum() + kpi["indirect_costs"].sum()) / kpi["earned_premium"].sum() * 100
        if kpi["earned_premium"].sum() else np.nan))
    c4.metric("Profitability", fmt_omr(kpi["profitability"].sum()))
    var = variance_analysis(df)
    c5.metric("Budget Variance", fmt_pct(var["variance_pct"]), fmt_omr(var["variance_abs"]))

    st.markdown("---")
    col1, col2 = st.columns(2)
    with col1:
        fig = px.line(growth, x="period", y="gross_written_premium", markers=True,
                       title="Gross Written Premium Trend")
        fig.update_traces(line_color="#1f4e79")
        st.plotly_chart(fig, use_container_width=True)
    with col2:
        by_group = kpi.groupby("group_of_product", as_index=False)["gross_written_premium"].sum()
        fig2 = px.pie(by_group, names="group_of_product", values="gross_written_premium",
                       title="Premium Mix by Line of Business", hole=0.45)
        st.plotly_chart(fig2, use_container_width=True)

    st.markdown("#### Recent Exceptions")
    outliers = detect_outliers(kpi, method=cfg.get("outlier_method", "Z-Score"),
                                threshold=cfg.get("outlier_threshold", 3.0))
    flagged = outliers[outliers["is_outlier"]].sort_values("outlier_score", ascending=False).head(5)
    if flagged.empty:
        st.success("No significant exceptions detected in the current filter.")
    else:
        st.dataframe(flagged[["year", "month", "branch", "product", "gross_written_premium", "outlier_score"]],
                     use_container_width=True, hide_index=True)

    kpi_glossary_expander()


# =============================================================================
# PAGE: SETUP
# =============================================================================

def page_setup(cfg: dict, hierarchy: pd.DataFrame):
    st.title("⚙️ Setup")
    st.caption("Define your product hierarchy and company configuration before uploading or analyzing data.")

    tab1, tab2 = st.tabs(["Product Hierarchy", "Company Settings"])

    with tab1:
        st.markdown("#### Group of Product → Product → Sub-product")
        st.caption("The Group of Product also serves as the Line of Business, so every LOB flows consistently "
                    "through filters, uploads, forecasts and reports.")
        edited = st.data_editor(
            hierarchy, num_rows="dynamic", use_container_width=True,
            column_config={
                "group_of_product": st.column_config.TextColumn("Group of Product / LOB", required=True),
                "product": st.column_config.TextColumn("Product", required=True),
                "sub_product": st.column_config.TextColumn("Sub-product", required=True),
            },
            key="hierarchy_editor",
        )
        col1, col2 = st.columns([1, 5])
        with col1:
            if st.button("💾 Save Hierarchy", type="primary"):
                errors = validate_hierarchy(edited)
                if errors:
                    for e in errors:
                        st.error(e)
                else:
                    save_hierarchy(edited)
                    st.success("Hierarchy saved. It will now flow into filters, uploads, forecasts and reports.")
                    st.rerun()

    with tab2:
        with st.form("company_settings_form"):
            company_name = st.text_input("Company Name", value=cfg["company_name"])
            currency_label = st.text_input("Reporting Currency Label", value=cfg["currency_label"])
            fy_start = st.selectbox("Fiscal Year Start Month", list(range(1, 13)),
                                     index=cfg["fiscal_year_start_month"] - 1)
            submitted = st.form_submit_button("Save Settings", type="primary")
            if submitted:
                cfg.update({
                    "company_name": company_name,
                    "currency_label": currency_label,
                    "fiscal_year_start_month": fy_start,
                })
                save_config(cfg)
                st.success("Company settings saved.")
                st.rerun()


# =============================================================================
# PAGE: DATA UPLOAD
# =============================================================================

def page_data_upload(hierarchy: pd.DataFrame):
    st.title("📤 Data Upload")
    if hierarchy.empty:
        st.error("Please define your Product Hierarchy in **Setup** before uploading data.")
        return

    template_type = st.selectbox(
        "Template type", ["Budget", "Actual", "Forecast", "Claims", "Production", "Expenses"]
    )
    template_df = make_template(template_type)

    col1, col2 = st.columns(2)
    with col1:
        st.download_button(
            f"⬇️ Download {template_type} Template (CSV)",
            data=export_csv(template_df),
            file_name=f"{template_type.lower()}_template.csv",
            mime="text/csv",
        )
    with col2:
        st.download_button(
            f"⬇️ Download {template_type} Template (Excel)",
            data=export_excel({template_type: template_df}),
            file_name=f"{template_type.lower()}_template.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        )

    st.markdown("---")
    uploaded = st.file_uploader("Upload completed file (CSV or Excel)", type=["csv", "xlsx"])
    if uploaded is not None:
        try:
            if uploaded.name.endswith(".csv"):
                new_df = pd.read_csv(uploaded)
            else:
                new_df = pd.read_excel(uploaded)
        except Exception as e:
            st.error(f"Could not read file: {e}")
            return

        st.markdown("#### Preview")
        st.dataframe(new_df.head(20), use_container_width=True)

        errors = validate_upload(new_df, template_type)
        if errors:
            st.error("Validation failed:")
            for e in errors:
                st.write(f"- {e}")
        else:
            st.success(f"✅ {len(new_df)} row(s) validated successfully.")
            if st.button("Confirm and Merge into Dataset", type="primary"):
                existing = load_data()
                merged = merge_uploaded_data(existing, new_df)
                save_data(merged)
                st.success("Data merged into the master dataset.")
                st.balloons()


# =============================================================================
# PAGE: CEO / COO / CFO VIEWS
# =============================================================================

def page_ceo_view(df: pd.DataFrame, cfg: dict):
    st.title("📈 CEO View — Strategic Performance")
    if df.empty:
        st.info("No data available.")
        return
    kpi = calculate_kpis(df)
    var = variance_analysis(df)

    c1, c2, c3, c4 = st.columns(4)
    c1.metric("Total GWP", fmt_omr(kpi["gross_written_premium"].sum()))
    c2.metric("Profitability", fmt_omr(kpi["profitability"].sum()))
    c3.metric("Combined Ratio", fmt_pct(kpi["combined_ratio"].mean()))
    c4.metric("Budget vs Actual", fmt_pct(var["variance_pct"]))

    growth = period_growth(kpi)
    fig = px.area(growth, x="period", y="gross_written_premium", title="Portfolio Growth Trend")
    st.plotly_chart(fig, use_container_width=True)

    col1, col2 = st.columns(2)
    with col1:
        by_group = kpi.groupby("group_of_product", as_index=False)["profitability"].sum()
        st.plotly_chart(px.bar(by_group, x="group_of_product", y="profitability",
                                title="Profitability by Line of Business", color="profitability",
                                color_continuous_scale="RdYlGn"), use_container_width=True)
    with col2:
        outliers = detect_outliers(kpi, method=cfg.get("outlier_method", "Z-Score"),
                                    threshold=cfg.get("outlier_threshold", 3.0))
        top_exceptions = outliers[outliers["is_outlier"]].nlargest(10, "outlier_score")
        st.markdown("**Top Exceptions**")
        if top_exceptions.empty:
            st.success("No major exceptions.")
        else:
            st.dataframe(top_exceptions[["branch", "product", "gross_written_premium", "outlier_score"]],
                         use_container_width=True, hide_index=True)

    st.markdown("**Budget vs Actual vs Forecast (by Line of Business)**")
    st.dataframe(var["by_group"], use_container_width=True, hide_index=True)

    render_ai_insights_and_export("CEO", df, cfg)


def page_coo_view(df: pd.DataFrame, cfg: dict):
    st.title("⚙️ COO View — Operations & Production")
    if df.empty:
        st.info("No data available.")
        return
    kpi = calculate_kpis(df)

    c1, c2, c3, c4 = st.columns(4)
    c1.metric("Claims Reported", fmt_omr(kpi["claims_reported"].sum()))
    c2.metric("Claims Paid", fmt_omr(kpi["claims_paid"].sum()))
    c3.metric("Outstanding Claims", fmt_omr(kpi["outstanding_claims"].sum()))
    c4.metric("Claims Ratio", fmt_pct(kpi["claims_ratio"].mean()))

    col1, col2 = st.columns(2)
    with col1:
        by_branch = kpi.groupby("branch", as_index=False)["gross_written_premium"].sum()
        st.plotly_chart(px.bar(by_branch, x="branch", y="gross_written_premium",
                                title="Production by Branch"), use_container_width=True)
    with col2:
        by_channel = kpi.groupby("channel", as_index=False)["gross_written_premium"].sum()
        st.plotly_chart(px.bar(by_channel, x="channel", y="gross_written_premium",
                                title="Production by Channel", color="channel"),
                         use_container_width=True)

    st.markdown("**Outlier Production (flagged)**")
    outliers = detect_outliers(kpi, group_cols=["branch", "agent"],
                                method=cfg.get("outlier_method", "Z-Score"),
                                threshold=cfg.get("outlier_threshold", 3.0))
    flagged = outliers[outliers["is_outlier"]]
    st.dataframe(flagged[["branch", "agent", "product", "gross_written_premium", "outlier_score"]]
                 if not flagged.empty else pd.DataFrame(columns=["No exceptions"]),
                 use_container_width=True, hide_index=True)

    render_ai_insights_and_export("COO", df, cfg)


def page_cfo_view(df: pd.DataFrame, cfg: dict):
    st.title("💰 CFO View — Financial Health")
    if df.empty:
        st.info("No data available.")
        return
    kpi = calculate_kpis(df)
    var = variance_analysis(df)

    c1, c2, c3, c4 = st.columns(4)
    c1.metric("Loss Ratio", fmt_pct(kpi["loss_ratio"].mean()))
    c2.metric("Expense Ratio", fmt_pct(kpi["expense_ratio"].mean()))
    c3.metric("Receivables", fmt_omr(kpi["receivables"].sum()))
    c4.metric("Collection Ratio", fmt_pct(kpi["collection_ratio"].mean()))

    col1, col2 = st.columns(2)
    with col1:
        trend = kpi.groupby(["year", "month"], as_index=False).agg(
            loss_ratio=("loss_ratio", "mean"), expense_ratio=("expense_ratio", "mean"),
            combined_ratio=("combined_ratio", "mean"))
        trend["period"] = pd.to_datetime(dict(year=trend.year, month=trend.month, day=1))
        fig = go.Figure()
        for col_name, label in [("loss_ratio", "Loss Ratio"), ("expense_ratio", "Expense Ratio"),
                                 ("combined_ratio", "Combined Ratio")]:
            fig.add_trace(go.Scatter(x=trend["period"], y=trend[col_name], mode="lines+markers", name=label))
        fig.update_layout(title="Ratio Trends", yaxis_title="%")
        st.plotly_chart(fig, use_container_width=True)
    with col2:
        forecast_df = forecast_metrics(kpi, horizon=cfg.get("forecast_horizon_months", 12))
        if not forecast_df.empty:
            fig2 = px.line(forecast_df, x="period", y="value", color="type",
                            title="Forecast Accuracy Outlook")
            st.plotly_chart(fig2, use_container_width=True)

    st.markdown("**Variance Analysis**")
    st.metric("Total Variance vs Budget", fmt_omr(var["variance_abs"]), fmt_pct(var["variance_pct"]))
    st.dataframe(var["by_group"], use_container_width=True, hide_index=True)
    kpi_glossary_expander()

    render_ai_insights_and_export("CFO", df, cfg)


# =============================================================================
# PAGE: OUTLIER DETECTION
# =============================================================================

def page_outliers(df: pd.DataFrame, cfg: dict):
    st.title("🔍 Outlier Detection")
    if df.empty:
        st.info("No data available.")
        return
    kpi = calculate_kpis(df)

    col1, col2, col3 = st.columns(3)
    method = col1.selectbox("Method", ["Z-Score", "IQR", "Isolation Forest"],
                             index=["Z-Score", "IQR", "Isolation Forest"].index(cfg.get("outlier_method", "Z-Score")))
    threshold = col2.slider("Z-Score Threshold", 1.5, 5.0, float(cfg.get("outlier_threshold", 3.0)), 0.1)
    group_by = col3.multiselect("Group by", ["branch", "agent", "channel", "product", "sub_product"],
                                 default=["branch"])
    value_col = st.selectbox("Metric", ["gross_written_premium", "claims_paid", "receivables", "actual"])

    result = detect_outliers(kpi, value_col=value_col, group_cols=group_by or ["branch"],
                              method=method, threshold=threshold)
    flagged = result[result["is_outlier"]].sort_values("outlier_score", ascending=False)

    c1, c2 = st.columns(2)
    c1.metric("Total Records", len(result))
    c2.metric("Flagged Exceptions", len(flagged))

    fig = px.scatter(result, x=result.index, y=value_col, color="is_outlier",
                      color_discrete_map={True: "#d62728", False: "#1f4e79"},
                      title=f"Exception Chart — {value_col}")
    st.plotly_chart(fig, use_container_width=True)

    st.markdown("#### Flagged Records")
    cols_to_show = ["year", "month", "branch", "agent", "channel", "product", "sub_product",
                     value_col, "outlier_score"]
    st.dataframe(flagged[cols_to_show] if not flagged.empty else pd.DataFrame(columns=cols_to_show),
                 use_container_width=True, hide_index=True)

    if not flagged.empty:
        st.download_button("⬇️ Export Flagged Records (CSV)", export_csv(flagged[cols_to_show]),
                            "outliers.csv", "text/csv")


# =============================================================================
# PAGE: FORECASTING
# =============================================================================

def page_forecasting(df: pd.DataFrame, cfg: dict):
    st.title("🔮 Auto-Forecasting")
    if df.empty:
        st.info("No data available.")
        return

    metric = st.selectbox("Metric to forecast", ["gross_written_premium", "claims_paid", "earned_premium"])
    horizon = st.slider("Forecast horizon (months)", 3, 24, cfg.get("forecast_horizon_months", 12))

    kpi = calculate_kpis(df)
    fdf = forecast_metrics(kpi, value_col=metric, horizon=horizon)
    if fdf.empty:
        st.warning("Not enough historical data to forecast (need at least 4 months).")
        return

    fig = go.Figure()
    hist = fdf[fdf["type"] == "Historical"]
    fut = fdf[fdf["type"] == "Forecast"]
    fig.add_trace(go.Scatter(x=hist["period"], y=hist["value"], name="Historical",
                              mode="lines+markers", line=dict(color="#1f4e79")))
    fig.add_trace(go.Scatter(x=fut["period"], y=fut["value"], name="Forecast",
                              mode="lines+markers", line=dict(color="#d62728", dash="dash")))
    fig.add_trace(go.Scatter(x=fut["period"], y=fut["upper"], name="Upper CI",
                              line=dict(width=0), showlegend=False))
    fig.add_trace(go.Scatter(x=fut["period"], y=fut["lower"], name="Confidence Interval",
                              fill="tonexty", line=dict(width=0), fillcolor="rgba(214,39,40,0.15)"))
    fig.update_layout(title=f"{metric.replace('_', ' ').title()} — Historical vs Forecast")
    st.plotly_chart(fig, use_container_width=True)

    st.markdown("#### Forecast Table")
    display_fut = fut[["period", "value", "lower", "upper"]].copy()
    for c in ["value", "lower", "upper"]:
        display_fut[c] = display_fut[c].round(0)
    st.dataframe(display_fut, use_container_width=True, hide_index=True)
    st.download_button("⬇️ Export Forecast (CSV)", export_csv(display_fut), "forecast.csv", "text/csv")


# =============================================================================
# PAGE: BUDGET VS ACTUAL VS FORECAST
# =============================================================================

def page_bvaf(df: pd.DataFrame):
    st.title("📊 Budget vs Actual vs Forecast")
    if df.empty:
        st.info("No data available.")
        return
    var = variance_analysis(df)

    c1, c2, c3 = st.columns(3)
    c1.metric("Budget", fmt_omr(var["total_budget"]))
    c2.metric("Actual", fmt_omr(var["total_actual"]), f"{var['variance_pct']:+.0f}%")
    c3.metric("Forecast", fmt_omr(var["total_forecast"]), f"{var['forecast_variance_pct']:+.0f}%")

    col1, col2 = st.columns(2)
    with col1:
        comp = pd.DataFrame({
            "Metric": ["Budget", "Actual", "Forecast"],
            "Value": [var["total_budget"], var["total_actual"], var["total_forecast"]],
        })
        st.plotly_chart(px.bar(comp, x="Metric", y="Value", color="Metric",
                                title="Budget vs Actual vs Forecast"), use_container_width=True)
    with col2:
        st.markdown("#### Variance Drivers (Waterfall)")
        drivers = var["drivers"]
        labels = ["Budget"] + list(drivers.keys()) + ["Actual"]
        values = [var["total_budget"]] + list(drivers.values()) + [0]
        measures = ["absolute"] + ["relative"] * len(drivers) + ["total"]
        fig = go.Figure(go.Waterfall(
            x=labels, measure=measures, y=values,
            connector={"line": {"color": "#999"}},
            decreasing={"marker": {"color": "#d62728"}},
            increasing={"marker": {"color": "#2ca02c"}},
            totals={"marker": {"color": "#1f4e79"}},
        ))
        fig.update_layout(title="Variance Driver Waterfall")
        st.plotly_chart(fig, use_container_width=True)

    st.markdown("#### Variance by Line of Business")
    by_group = var["by_group"].copy()
    for c in ["budget", "actual", "forecast", "variance_abs"]:
        by_group[c] = by_group[c].apply(fmt_omr)
    by_group["variance_pct"] = var["by_group"]["variance_pct"].apply(fmt_pct)
    st.dataframe(by_group, use_container_width=True, hide_index=True)

    st.markdown("#### Executive Summary")
    direction = "favorable" if var["variance_pct"] >= 0 else "unfavorable"
    st.info(
        f"Actual performance shows a **{direction}** variance of {fmt_omr(var['variance_abs'])} "
        f"({fmt_pct(var['variance_pct'])}) against budget, driven primarily by "
        f"{'production growth' if var['drivers']['Volume/Production'] >= 0 else 'production shortfall'} "
        f"and claims/expense movements."
    )


# =============================================================================
# PAGE: AI INSIGHTS
# =============================================================================

def page_ai_insights(df: pd.DataFrame, cfg: dict):
    st.title("🤖 AI Insights")
    if df.empty:
        st.info("No data available.")
        return
    kpi = calculate_kpis(df)
    var = variance_analysis(df)
    outliers = detect_outliers(kpi, method=cfg.get("outlier_method", "Z-Score"),
                                threshold=cfg.get("outlier_threshold", 3.0))
    outliers_count = int(outliers["is_outlier"].sum())

    kpi_summary = {
        "loss_ratio": kpi["loss_ratio"].mean(),
        "expense_ratio": kpi["expense_ratio"].mean(),
        "combined_ratio": kpi["combined_ratio"].mean(),
        "collection_ratio": kpi["collection_ratio"].mean(),
        "profitability": kpi["profitability"].sum(),
    }

    if st.button("🔄 Regenerate Insights", type="primary"):
        st.session_state.pop("_ai_insights_cache", None)

    if "_ai_insights_cache" not in st.session_state:
        with st.spinner("Analyzing KPIs, outliers, and variances..."):
            st.session_state["_ai_insights_cache"] = generate_ai_insights(
                kpi_summary, outliers_count, var,
                api_key=cfg.get("gemini_api_key", ""),
                model_name=cfg.get("gemini_model", GEMINI_MODEL))

    st.markdown(st.session_state["_ai_insights_cache"])

    if not cfg.get("gemini_api_key"):
        st.caption(
            "ℹ️ Running in rule-based mode. Add a Gemini API key in **Admin / Configuration** "
            f"to enable AI-generated narrative commentary (model: `{cfg.get('gemini_model', GEMINI_MODEL)}`)."
        )
    else:
        st.caption(f"ℹ️ Powered by `{cfg.get('gemini_model', GEMINI_MODEL)}`.")


# =============================================================================
# PAGE: REPORTS
# =============================================================================

def page_reports(df: pd.DataFrame, cfg: dict):
    st.title("📄 Reports")
    if df.empty:
        st.info("No data available.")
        return
    var = variance_analysis(df)
    insights_text = st.session_state.get("_ai_insights_cache", "Run AI Insights first for narrative commentary.")

    tab1, tab2 = st.tabs(["📘 Full Management Report", "Standard Reports"])

    with tab1:
        st.markdown(
            "A single rich, board-ready PDF covering: Executive Summary (AI Insights), a full "
            "**KPI Dashboard with definitions**, detailed **Variance Analysis**, and the "
            "**Production Forecast** — chart and table."
        )
        if st.button("📘 Generate Full Management Report (PDF)", type="primary"):
            with st.spinner("Building the full management report..."):
                content = generate_management_report(df, cfg, insights_text)
            st.success("Management report generated.")
            st.download_button(
                "⬇️ Download Management Report (PDF)", data=content,
                file_name="smart_analytics_management_report.pdf", mime="application/pdf",
            )

    with tab2:
        report_type = st.selectbox("Report format", ["PDF", "Excel", "PPTX", "DOCX"])
        if report_type == "PDF":
            st.caption("A concise portfolio summary PDF that includes headline financials and AI Executive Insights.")
        if st.button("Generate Report", type="primary"):
            try:
                with st.spinner(f"Generating {report_type} report..."):
                    content = generate_reports(report_type, df, var, insights_text, cfg["company_name"])
                mime_map = {
                    "PDF": "application/pdf",
                    "Excel": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    "PPTX": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    "DOCX": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                }
                ext_map = {"PDF": "pdf", "Excel": "xlsx", "PPTX": "pptx", "DOCX": "docx"}
                st.success("Report generated.")
                st.download_button(
                    f"⬇️ Download {report_type} Report", data=content,
                    file_name=f"smart_analytics_portfolio_report.{ext_map[report_type]}",
                    mime=mime_map[report_type],
                )
            except ImportError as e:
                st.error(
                    f"Missing package for {report_type} generation: {e}. "
                    f"See SETUP_GUIDE.md for the required `pip install` command."
                )

    st.markdown("---")
    st.markdown("#### Quick Export")
    col1, col2 = st.columns(2)
    col1.download_button("⬇️ Export Raw Data (CSV)", export_csv(df), "raw_data.csv", "text/csv")
    col2.download_button("⬇️ Export Raw Data (Excel)", export_excel({"Data": df}),
                          "raw_data.xlsx",
                          "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")


# =============================================================================
# PAGE: ADMIN / CONFIGURATION
# =============================================================================

def page_admin(cfg: dict, hierarchy: pd.DataFrame):
    st.title("🛠️ Admin / Configuration")
    tab1, tab2, tab3, tab4 = st.tabs(["AI & Outlier Settings", "User Management", "Sample Data", "Support Outbox"])

    with tab1:
        with st.form("ai_settings_form"):
            api_key = st.text_input("Gemini API Key", value=cfg.get("gemini_api_key", ""), type="password")
            model_name = st.text_input("Gemini Model", value=cfg.get("gemini_model", GEMINI_MODEL))
            method = st.selectbox("Default Outlier Method", ["Z-Score", "IQR", "Isolation Forest"],
                                   index=["Z-Score", "IQR", "Isolation Forest"].index(cfg.get("outlier_method", "Z-Score")))
            threshold = st.slider("Default Z-Score Threshold", 1.5, 5.0, float(cfg.get("outlier_threshold", 3.0)))
            horizon = st.slider("Default Forecast Horizon (months)", 3, 24, cfg.get("forecast_horizon_months", 12))
            if st.form_submit_button("Save", type="primary"):
                cfg.update({"gemini_api_key": api_key, "gemini_model": model_name, "outlier_method": method,
                            "outlier_threshold": threshold, "forecast_horizon_months": horizon})
                save_config(cfg)
                st.success("Configuration saved.")

    with tab2:
        users = load_users()
        st.dataframe(pd.DataFrame([
            {"username": u, "name": v["name"], "role": v["role"], "email": v["email"]}
            for u, v in users.items()
        ]), use_container_width=True, hide_index=True)

        with st.expander("Add / Update User"):
            with st.form("user_form"):
                uname = st.text_input("Username")
                name = st.text_input("Full Name")
                role = st.selectbox("Role", ["Admin", "CEO", "COO", "CFO"])
                email = st.text_input("Email")
                pw = st.text_input("Password", type="password")
                if st.form_submit_button("Save User", type="primary"):
                    if uname and pw:
                        users[uname] = {"name": name, "role": role, "email": email,
                                         "password_hash": _hash_pw(pw)}
                        save_users(users)
                        st.success(f"User '{uname}' saved.")
                        st.rerun()
                    else:
                        st.error("Username and password are required.")

    with tab3:
        st.write("Generate realistic sample data — covering every Line of Business in your product hierarchy — "
                 "for testing.")
        n_months = st.slider("Months of history", 6, 60, 24)
        if st.button("Generate Sample Data", type="primary"):
            sample = generate_sample_data(hierarchy, n_months=n_months)
            save_data(sample)
            st.success(f"Generated {len(sample)} sample rows across {n_months} months and "
                       f"{hierarchy['group_of_product'].nunique()} Lines of Business.")
            st.rerun()
        if st.button("Clear All Data", type="secondary"):
            save_data(empty_dataset())
            st.warning("All transaction data cleared.")
            st.rerun()

    with tab4:
        outbox = st.session_state.get("_email_outbox", [])
        if outbox:
            st.dataframe(pd.DataFrame(outbox), use_container_width=True, hide_index=True)
        else:
            st.caption("No support emails sent this session.")


# =============================================================================
# MAIN ROUTER
# =============================================================================

def main():
    if not st.session_state["authenticated"]:
        login_screen()
        return

    cfg = load_config()
    hierarchy = load_hierarchy()
    page, filters = sidebar_nav(cfg, hierarchy)
    raw_df = load_data()
    df = apply_filters(raw_df, filters) if not raw_df.empty else raw_df

    role = st.session_state["user"]["role"]
    restricted = {
        "CEO View": {"CEO", "Admin"},
        "COO View": {"COO", "Admin"},
        "CFO View": {"CFO", "Admin"},
        "Admin / Configuration": {"Admin"},
    }
    if page in restricted and role not in restricted[page]:
        st.error(f"Your role ({role}) does not have access to this view.")
        return

    if page == "Home Dashboard":
        page_home(df, cfg)
    elif page == "Setup":
        page_setup(cfg, hierarchy)
    elif page == "Data Upload":
        page_data_upload(hierarchy)
    elif page == "CEO View":
        page_ceo_view(df, cfg)
    elif page == "COO View":
        page_coo_view(df, cfg)
    elif page == "CFO View":
        page_cfo_view(df, cfg)
    elif page == "Outlier Detection":
        page_outliers(df, cfg)
    elif page == "Forecasting":
        page_forecasting(df, cfg)
    elif page == "Budget vs Actual vs Forecast":
        page_bvaf(df)
    elif page == "AI Insights":
        page_ai_insights(df, cfg)
    elif page == "Reports":
        page_reports(df, cfg)
    elif page == "Admin / Configuration":
        page_admin(cfg, hierarchy)


if __name__ == "__main__":
    main()
